from __future__ import annotations

import hashlib
import json
import mimetypes
import os
import shlex
import shutil
import subprocess
from dataclasses import dataclass
from datetime import timedelta
from pathlib import Path
from typing import Any

from sqlmodel import Session, select

from rag.common import append_jsonl, chunk_text, sanitize_metadata, write_json

from .config import OrquestraSettings
from .gateway import OrquestraGateway
from .models import (
    MemoryRecord,
    WorkspaceAsset,
    WorkspaceDerivative,
    WorkspaceInsight,
    WorkspaceScan,
    utc_now,
)
from .vector_index import OrquestraVectorIndex, blend_scores, score_overlap

try:  # pragma: no cover - depende do ambiente local
    from PIL import Image
except Exception:  # pragma: no cover - fallback sem pillow
    Image = None

try:  # pragma: no cover - depende do ambiente local
    from pypdf import PdfReader
except Exception:  # pragma: no cover
    PdfReader = None

try:  # pragma: no cover - depende do ambiente local
    from docx import Document as DocxDocument
except Exception:  # pragma: no cover
    DocxDocument = None

try:  # pragma: no cover - depende do ambiente local
    from openpyxl import load_workbook
except Exception:  # pragma: no cover
    load_workbook = None

try:  # pragma: no cover - depende do ambiente local
    from pptx import Presentation
except Exception:  # pragma: no cover
    Presentation = None

try:  # pragma: no cover - depende do ambiente local
    import mutagen
except Exception:  # pragma: no cover
    mutagen = None


TEXT_EXTENSIONS = {
    ".md",
    ".txt",
    ".json",
    ".jsonl",
    ".py",
    ".sh",
    ".yaml",
    ".yml",
    ".toml",
    ".csv",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".html",
    ".css",
    ".scss",
    ".sql",
    ".swift",
    ".java",
    ".kt",
    ".go",
    ".rs",
}
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff", ".heic"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".aac", ".m4a", ".ogg", ".flac"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".m4v", ".webm", ".mkv", ".avi"}
PDF_EXTENSIONS = {".pdf"}


@dataclass
class WorkspacePaths:
    root: Path
    inventory_dir: Path
    previews_dir: Path
    derivatives_dir: Path
    insights_dir: Path

    @classmethod
    def from_settings(cls, settings: OrquestraSettings) -> "WorkspacePaths":
        root = settings.artifacts_root / "workspace"
        return cls(
            root=root,
            inventory_dir=root / "inventories",
            previews_dir=root / "previews",
            derivatives_dir=root / "derivatives",
            insights_dir=root / "insights",
        )

    def ensure(self) -> None:
        for path in (self.root, self.inventory_dir, self.previews_dir, self.derivatives_dir, self.insights_dir):
            path.mkdir(parents=True, exist_ok=True)


def _slugify(raw: str) -> str:
    import re

    normalized = re.sub(r"[^a-zA-Z0-9]+", "-", raw).strip("-").lower()
    return normalized or "item"


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        while True:
            chunk = handle.read(1024 * 1024)
            if not chunk:
                break
            digest.update(chunk)
    return digest.hexdigest()


def _guess_kind(path: Path, mime_type: str) -> str:
    suffix = path.suffix.lower()
    if suffix in TEXT_EXTENSIONS or mime_type.startswith("text/"):
        return "code_text"
    if suffix in IMAGE_EXTENSIONS or mime_type.startswith("image/"):
        return "image"
    if suffix in PDF_EXTENSIONS or mime_type == "application/pdf":
        return "pdf"
    if suffix in OFFICE_EXTENSIONS:
        return "office"
    if suffix in AUDIO_EXTENSIONS or mime_type.startswith("audio/"):
        return "audio"
    if suffix in VIDEO_EXTENSIONS or mime_type.startswith("video/"):
        return "video"
    return "binary"


def _excerpt(text: str, size: int = 400) -> str:
    normalized = " ".join(text.split())
    if len(normalized) <= size:
        return normalized
    return normalized[: max(size - 1, 1)].rstrip() + "…"


class WorkspaceService:
    def __init__(self, settings: OrquestraSettings) -> None:
        self.settings = settings
        self.paths = WorkspacePaths.from_settings(settings)
        self.paths.ensure()
        self.index = OrquestraVectorIndex(settings)

    def attach_directory(
        self,
        session: Session,
        *,
        root_path: str,
        project_id: str | None,
        prompt_hint: str = "",
    ) -> WorkspaceScan:
        root = Path(root_path).expanduser().resolve()
        if not root.exists() or not root.is_dir():
            raise FileNotFoundError(f"Diretório inválido: {root}")

        scan = WorkspaceScan(
            project_id=project_id,
            root_path=str(root),
            prompt_hint=prompt_hint,
            status="scanning",
            inventory_path="",
        )
        session.add(scan)
        session.flush()

        inventory_path = self.paths.inventory_dir / f"{_slugify(scan.id)}.jsonl"
        total_assets = 0
        total_bytes = 0
        for file_path in sorted(root.rglob("*")):
            if not file_path.is_file():
                continue
            relative = file_path.relative_to(root)
            mime_type, _ = mimetypes.guess_type(str(file_path))
            mime_type = mime_type or "application/octet-stream"
            stat = file_path.stat()
            absolute = file_path.resolve()
            total_assets += 1
            total_bytes += stat.st_size
            asset = WorkspaceAsset(
                scan_id=scan.id,
                absolute_path=str(absolute),
                relative_path=str(relative),
                parent_relative_path=str(relative.parent) if str(relative.parent) != "." else None,
                asset_kind=_guess_kind(file_path, mime_type),
                mime_type=mime_type,
                extension=file_path.suffix.lower(),
                size_bytes=stat.st_size,
                sha256=_sha256(file_path),
                depth=max(len(relative.parts) - 1, 0),
                modified_at=utc_now(),
                title=file_path.name,
                summary_excerpt="",
                extraction_state="inventory_only",
                metadata_json=json.dumps(
                    {
                        "mtime_ns": stat.st_mtime_ns,
                        "prompt_hint": prompt_hint,
                    },
                    ensure_ascii=False,
                ),
            )
            session.add(asset)
            session.flush()
            append_jsonl(
                inventory_path,
                {
                    "asset_id": asset.id,
                    "absolute_path": asset.absolute_path,
                    "relative_path": asset.relative_path,
                    "asset_kind": asset.asset_kind,
                    "mime_type": asset.mime_type,
                    "size_bytes": asset.size_bytes,
                    "sha256": asset.sha256,
                    "depth": asset.depth,
                },
            )

        scan.inventory_path = str(inventory_path)
        scan.total_assets = total_assets
        scan.total_bytes = total_bytes
        scan.status = "ready"
        scan.updated_at = utc_now()
        session.add(
            WorkspaceInsight(
                scan_id=scan.id,
                kind="inventory_summary",
                title="Inventário inicial concluído",
                content=f"Varredura concluída em {root} com {total_assets} arquivos e {total_bytes} bytes.",
                relevance=0.4,
                metadata_json=json.dumps({"root_path": str(root)}, ensure_ascii=False),
            )
        )
        return scan

    def extract_asset(
        self,
        session: Session,
        asset: WorkspaceAsset,
        *,
        force: bool = False,
        prompt_hint: str = "",
    ) -> dict[str, Any]:
        if asset.extraction_state == "ready" and not force:
            return self.preview_asset(session, asset)

        handler = {
            "code_text": self._extract_code_text,
            "image": self._extract_image,
            "pdf": self._extract_pdf,
            "office": self._extract_office,
            "audio": self._extract_audio,
            "video": self._extract_video,
            "binary": self._extract_binary,
        }.get(asset.asset_kind, self._extract_binary)
        result = handler(asset, prompt_hint=prompt_hint)
        asset.summary_excerpt = result.get("summary_excerpt", asset.summary_excerpt)
        asset.extraction_state = result.get("state", "ready")
        metadata = json.loads(asset.metadata_json or "{}")
        metadata.update(result.get("metadata", {}))
        asset.metadata_json = json.dumps(metadata, ensure_ascii=False)
        asset.updated_at = utc_now()
        session.add(asset)
        self._sync_derivatives(session, asset, result.get("derivatives", []))
        self._index_asset(asset, result)
        return self.preview_asset(session, asset)

    def preview_asset(self, session: Session, asset: WorkspaceAsset) -> dict[str, Any]:
        derivatives = session.exec(select(WorkspaceDerivative).where(WorkspaceDerivative.asset_id == asset.id)).all()
        metadata = json.loads(asset.metadata_json or "{}")
        raw_url = f"/api/workspace/assets/{asset.id}/preview?raw=true"
        preview_type = asset.asset_kind if asset.asset_kind != "binary" else "metadata"
        return {
            "asset_id": asset.id,
            "title": asset.title,
            "asset_kind": asset.asset_kind,
            "mime_type": asset.mime_type,
            "summary_excerpt": asset.summary_excerpt,
            "metadata": metadata,
            "preview_type": preview_type,
            "raw_url": raw_url,
            "derivatives": [
                {
                    "id": item.id,
                    "kind": item.derivative_kind,
                    "media_type": item.media_type,
                    "storage_path": item.storage_path,
                }
                for item in derivatives
            ],
        }

    def query_workspace(
        self,
        session: Session,
        gateway: OrquestraGateway,
        *,
        scan_id: str,
        prompt: str,
        provider_id: str | None = None,
        model_name: str | None = None,
        force_extract: bool = False,
    ) -> dict[str, Any]:
        scan = session.get(WorkspaceScan, scan_id)
        if scan is None:
            raise FileNotFoundError(f"Scan não encontrado: {scan_id}")
        assets = session.exec(select(WorkspaceAsset).where(WorkspaceAsset.scan_id == scan_id)).all()
        ranked = self._rank_assets(prompt, assets)
        selected = ranked[:8]
        extracted: list[dict[str, Any]] = []
        for item in selected:
            asset = next(asset for asset in assets if asset.id == item["asset_id"])
            if asset.extraction_state != "ready" or force_extract:
                self.extract_asset(session, asset, force=force_extract, prompt_hint=prompt)
            extracted.append(
                {
                    "asset_id": asset.id,
                    "title": asset.title,
                    "relative_path": asset.relative_path,
                    "asset_kind": asset.asset_kind,
                    "summary_excerpt": asset.summary_excerpt,
                    "score": item["score"],
                    "metadata": json.loads(asset.metadata_json or "{}"),
                }
            )

        context_parts = []
        for item in extracted:
            text = item["summary_excerpt"] or ""
            context_parts.append(f"[{item['asset_kind'].upper()}] {item['relative_path']}\n{text}")
        context = "\n\n".join(part for part in context_parts if part.strip())
        fallback = context_parts[0] if context_parts else "Nenhum arquivo relevante encontrado."
        response = gateway.generate(
            messages=[
                {
                    "role": "system",
                    "content": (
                        "Você é o assistente do Workspace Multimodal do Orquestra. "
                        "Analise o inventário anexado, cite caminhos de arquivo relevantes, "
                        "explique como abrir ou personalizar os arquivos e seja objetivo."
                    ),
                },
                {
                    "role": "user",
                    "content": f"Inventário anexado: {scan.root_path}\n\nContexto:\n{context}\n\nPergunta:\n{prompt}",
                },
            ],
            provider_id=provider_id,
            model_name=model_name,
            fallback_text=fallback,
        )
        insight = WorkspaceInsight(
            scan_id=scan.id,
            kind="query_response",
            title="Consulta multimodal",
            content=response.content,
            relevance=0.82,
            metadata_json=json.dumps(
                {
                    "provider_id": response.provider_id,
                    "model_name": response.model_name,
                    "asset_ids": [item["asset_id"] for item in extracted],
                },
                ensure_ascii=False,
            ),
        )
        session.add(insight)
        return {
            "scan_id": scan.id,
            "answer": response.content,
            "provider_id": response.provider_id,
            "model_name": response.model_name,
            "usage": response.usage,
            "latency_seconds": response.latency_seconds,
            "assets": extracted,
        }

    def memorize_asset(
        self,
        session: Session,
        asset: WorkspaceAsset,
        *,
        project_id: str | None,
        scope: str = "workspace_memory",
        source: str = "workspace",
    ) -> MemoryRecord:
        metadata = json.loads(asset.metadata_json or "{}")
        record = MemoryRecord(
            project_id=project_id,
            scope=scope,
            memory_kind="reference",
            source=f"{source}:{asset.relative_path}",
            content=f"Arquivo {asset.relative_path} ({asset.asset_kind}): {asset.summary_excerpt or asset.title}",
            confidence=0.69,
            metadata_json=json.dumps(metadata, ensure_ascii=False),
        )
        session.add(record)
        return record

    def open_asset(self, asset: WorkspaceAsset) -> dict[str, Any]:
        target = Path(asset.absolute_path)
        if not target.exists():
            raise FileNotFoundError(f"Arquivo não encontrado: {target}")
        if shutil.which("open"):
            command = ["open", str(target)]
        elif shutil.which("xdg-open"):
            command = ["xdg-open", str(target)]
        else:
            command = [str(target)]
        subprocess.Popen(command)
        return {"ok": True, "command": command, "path": str(target)}

    def gc_derivatives(self, session: Session) -> int:
        deleted = 0
        now = utc_now()
        rows = session.exec(select(WorkspaceDerivative)).all()
        for row in rows:
            if row.expires_at and row.expires_at <= now:
                path = Path(row.storage_path)
                if path.exists():
                    path.unlink()
                session.delete(row)
                deleted += 1
        return deleted

    def _index_asset(self, asset: WorkspaceAsset, result: dict[str, Any]) -> None:
        text = result.get("index_text") or result.get("summary_excerpt") or ""
        if not text.strip():
            return
        parts = chunk_text(text, chunk_size=900, overlap=120)
        self.index.upsert(
            "workspace_assets",
            [
                {
                    "id": f"{asset.id}:{index}",
                    "text": chunk,
                    "payload": {
                        "asset_id": asset.id,
                        "relative_path": asset.relative_path,
                        "asset_kind": asset.asset_kind,
                        "title": asset.title,
                    },
                }
                for index, chunk in enumerate(parts)
            ],
        )

    def _sync_derivatives(self, session: Session, asset: WorkspaceAsset, derivatives: list[dict[str, Any]]) -> None:
        for item in derivatives:
            derivative = WorkspaceDerivative(
                asset_id=asset.id,
                derivative_kind=item["kind"],
                storage_path=item["path"],
                media_type=item.get("media_type", ""),
                expires_at=utc_now() + timedelta(days=item.get("ttl_days", 7)),
                metadata_json=json.dumps(item.get("metadata", {}), ensure_ascii=False),
            )
            session.add(derivative)

    def _rank_assets(self, prompt: str, assets: list[WorkspaceAsset]) -> list[dict[str, Any]]:
        semantic_hits = self.index.query("workspace_assets", prompt, limit=12) if prompt.strip() else []
        semantic_map = {hit.payload.get("asset_id"): hit.score for hit in semantic_hits if hit.payload.get("asset_id")}
        ranked: list[dict[str, Any]] = []
        for asset in assets:
            lexical = score_overlap(prompt, asset.relative_path, asset.title, asset.summary_excerpt, asset.asset_kind)
            score = blend_scores(lexical, semantic_map.get(asset.id, 0.0))
            if prompt.strip():
                lowered = prompt.lower()
                if "imagem" in lowered and asset.asset_kind == "image":
                    score += 0.18
                if "pdf" in lowered and asset.asset_kind == "pdf":
                    score += 0.18
                if any(token in lowered for token in ("video", "vídeo")) and asset.asset_kind == "video":
                    score += 0.18
                if "audio" in lowered and asset.asset_kind == "audio":
                    score += 0.18
            else:
                score += max(0.02, 1 / max(asset.depth + 1, 1))
            ranked.append({"asset_id": asset.id, "score": round(score, 4)})
        return sorted(ranked, key=lambda item: item["score"], reverse=True)

    def _derivative_path(self, asset: WorkspaceAsset, suffix: str) -> Path:
        scan_dir = self.paths.derivatives_dir / _slugify(asset.scan_id)
        scan_dir.mkdir(parents=True, exist_ok=True)
        return scan_dir / f"{_slugify(asset.id)}{suffix}"

    def _extract_code_text(self, asset: WorkspaceAsset, *, prompt_hint: str = "") -> dict[str, Any]:
        text = Path(asset.absolute_path).read_text(encoding="utf-8", errors="ignore")
        summary = _excerpt(text, 700)
        return {
            "state": "ready",
            "summary_excerpt": summary,
            "index_text": text[:20000],
            "metadata": {
                "line_count": len(text.splitlines()),
                "prompt_hint": prompt_hint,
            },
        }

    def _extract_image(self, asset: WorkspaceAsset, *, prompt_hint: str = "") -> dict[str, Any]:
        metadata: dict[str, Any] = {"prompt_hint": prompt_hint}
        derivatives: list[dict[str, Any]] = []
        if Image is not None:
            with Image.open(asset.absolute_path) as image:
                metadata.update({"width": image.width, "height": image.height, "mode": image.mode})
                preview_path = self._derivative_path(asset, ".png")
                preview_image = image.copy()
                preview_image.thumbnail((640, 640))
                preview_image.save(preview_path, format="PNG")
                derivatives.append({"kind": "thumbnail", "path": str(preview_path), "media_type": "image/png", "ttl_days": 14})
        summary = f"Imagem {asset.title} com metadados {metadata}" if metadata else f"Imagem {asset.title}"
        return {"state": "ready", "summary_excerpt": _excerpt(summary, 320), "metadata": metadata, "derivatives": derivatives, "index_text": summary}

    def _extract_pdf(self, asset: WorkspaceAsset, *, prompt_hint: str = "") -> dict[str, Any]:
        metadata: dict[str, Any] = {"prompt_hint": prompt_hint}
        text = ""
        if PdfReader is not None:
            reader = PdfReader(asset.absolute_path)
            metadata["pages"] = len(reader.pages)
            text = "\n".join((page.extract_text() or "") for page in reader.pages[: min(len(reader.pages), 12)])
        summary = _excerpt(text or f"PDF {asset.title}", 900)
        derivative_path = self._derivative_path(asset, ".txt")
        derivative_path.write_text(text or summary, encoding="utf-8")
        return {
            "state": "ready",
            "summary_excerpt": summary,
            "metadata": metadata,
            "derivatives": [{"kind": "text_extract", "path": str(derivative_path), "media_type": "text/plain", "ttl_days": 21}],
            "index_text": text or summary,
        }

    def _extract_office(self, asset: WorkspaceAsset, *, prompt_hint: str = "") -> dict[str, Any]:
        suffix = asset.extension
        text = ""
        metadata: dict[str, Any] = {"prompt_hint": prompt_hint}
        if suffix == ".docx" and DocxDocument is not None:
            document = DocxDocument(asset.absolute_path)
            text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
        elif suffix == ".xlsx" and load_workbook is not None:
            workbook = load_workbook(asset.absolute_path, read_only=True, data_only=True)
            metadata["sheets"] = workbook.sheetnames
            lines: list[str] = []
            for sheet_name in workbook.sheetnames[:4]:
                sheet = workbook[sheet_name]
                lines.append(f"[Sheet] {sheet_name}")
                for row in sheet.iter_rows(min_row=1, max_row=20, values_only=True):
                    values = [str(cell) for cell in row if cell is not None]
                    if values:
                        lines.append(" | ".join(values))
            text = "\n".join(lines)
        elif suffix == ".pptx" and Presentation is not None:
            presentation = Presentation(asset.absolute_path)
            metadata["slides"] = len(presentation.slides)
            lines = []
            for slide_index, slide in enumerate(presentation.slides[:12], start=1):
                lines.append(f"[Slide {slide_index}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            text = "\n".join(lines)
        summary = _excerpt(text or f"Arquivo Office {asset.title}", 900)
        derivative_path = self._derivative_path(asset, ".txt")
        derivative_path.write_text(text or summary, encoding="utf-8")
        return {
            "state": "ready",
            "summary_excerpt": summary,
            "metadata": metadata,
            "derivatives": [{"kind": "text_extract", "path": str(derivative_path), "media_type": "text/plain", "ttl_days": 21}],
            "index_text": text or summary,
        }

    def _extract_audio(self, asset: WorkspaceAsset, *, prompt_hint: str = "") -> dict[str, Any]:
        metadata: dict[str, Any] = {"prompt_hint": prompt_hint}
        transcript_text = ""
        if mutagen is not None:
            try:
                audio = mutagen.File(asset.absolute_path)
                if audio is not None and getattr(audio, "info", None) is not None:
                    metadata["duration_seconds"] = round(float(audio.info.length), 2)
            except Exception:
                pass
        transcript_path = self._maybe_transcribe_audio(asset)
        derivatives: list[dict[str, Any]] = []
        if transcript_path and transcript_path.exists():
            transcript_text = transcript_path.read_text(encoding="utf-8", errors="ignore")
            derivatives.append({"kind": "transcript", "path": str(transcript_path), "media_type": "text/plain", "ttl_days": 14})
        summary = _excerpt(transcript_text or f"Áudio {asset.title} com metadados {metadata}", 900)
        return {
            "state": "ready" if transcript_text else "metadata_only",
            "summary_excerpt": summary,
            "metadata": metadata | {"transcript_available": bool(transcript_text)},
            "derivatives": derivatives,
            "index_text": transcript_text or summary,
        }

    def _extract_video(self, asset: WorkspaceAsset, *, prompt_hint: str = "") -> dict[str, Any]:
        metadata: dict[str, Any] = {"prompt_hint": prompt_hint}
        derivatives: list[dict[str, Any]] = []
        ffprobe = shutil.which("ffprobe")
        ffmpeg = shutil.which("ffmpeg")
        if ffprobe:
            try:
                result = subprocess.run(
                    [
                        ffprobe,
                        "-v",
                        "error",
                        "-show_entries",
                        "format=duration,size",
                        "-of",
                        "json",
                        asset.absolute_path,
                    ],
                    capture_output=True,
                    text=True,
                    check=True,
                    timeout=20,
                )
                payload = json.loads(result.stdout or "{}")
                metadata.update(payload.get("format", {}))
            except Exception:
                pass
        if ffmpeg:
            preview_path = self._derivative_path(asset, ".jpg")
            try:
                subprocess.run(
                    [
                        ffmpeg,
                        "-y",
                        "-i",
                        asset.absolute_path,
                        "-frames:v",
                        "1",
                        str(preview_path),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=40,
                )
                if preview_path.exists():
                    derivatives.append({"kind": "poster_frame", "path": str(preview_path), "media_type": "image/jpeg", "ttl_days": 14})
            except Exception:
                pass
        transcript_path = self._maybe_transcribe_video(asset)
        if transcript_path and transcript_path.exists():
            derivatives.append({"kind": "transcript", "path": str(transcript_path), "media_type": "text/plain", "ttl_days": 14})
            transcript_text = transcript_path.read_text(encoding="utf-8", errors="ignore")
        else:
            transcript_text = ""
        summary = _excerpt(transcript_text or f"Vídeo {asset.title} com metadados {metadata}", 900)
        return {
            "state": "ready" if transcript_text or derivatives else "metadata_only",
            "summary_excerpt": summary,
            "metadata": metadata | {"transcript_available": bool(transcript_text)},
            "derivatives": derivatives,
            "index_text": transcript_text or summary,
        }

    def _extract_binary(self, asset: WorkspaceAsset, *, prompt_hint: str = "") -> dict[str, Any]:
        summary = f"Arquivo binário {asset.title} ({asset.mime_type}) sem extrator específico."
        return {"state": "metadata_only", "summary_excerpt": summary, "metadata": {"prompt_hint": prompt_hint}, "index_text": summary}

    def _maybe_transcribe_audio(self, asset: WorkspaceAsset) -> Path | None:
        whisper = shutil.which("whisper")
        if whisper is None:
            return None
        output_dir = self._derivative_path(asset, "").parent
        output_dir.mkdir(parents=True, exist_ok=True)
        stem = self._derivative_path(asset, "").stem
        transcript_path = output_dir / f"{stem}.txt"
        if transcript_path.exists():
            return transcript_path
        try:
            subprocess.run(
                [
                    whisper,
                    asset.absolute_path,
                    "--model",
                    os.getenv("ORQUESTRA_WHISPER_MODEL", "turbo"),
                    "--output_dir",
                    str(output_dir),
                    "--output_format",
                    "txt",
                ],
                capture_output=True,
                text=True,
                timeout=300,
            )
        except Exception:
            return None
        return transcript_path if transcript_path.exists() else None

    def _maybe_transcribe_video(self, asset: WorkspaceAsset) -> Path | None:
        ffmpeg = shutil.which("ffmpeg")
        if ffmpeg is None:
            return None
        audio_path = self._derivative_path(asset, ".wav")
        if not audio_path.exists():
            try:
                subprocess.run(
                    [ffmpeg, "-y", "-i", asset.absolute_path, "-vn", str(audio_path)],
                    capture_output=True,
                    text=True,
                    timeout=300,
                )
            except Exception:
                return None
        fake_audio = WorkspaceAsset.model_construct(
            id=asset.id,
            scan_id=asset.scan_id,
            absolute_path=str(audio_path),
            relative_path=asset.relative_path,
            parent_relative_path=asset.parent_relative_path,
            asset_kind="audio",
            mime_type="audio/wav",
            extension=".wav",
            size_bytes=audio_path.stat().st_size if audio_path.exists() else 0,
            sha256="",
            depth=asset.depth,
            modified_at=utc_now(),
            title=asset.title,
            summary_excerpt="",
            extraction_state="inventory_only",
            metadata_json="{}",
            created_at=utc_now(),
            updated_at=utc_now(),
        )
        return self._maybe_transcribe_audio(fake_audio)
